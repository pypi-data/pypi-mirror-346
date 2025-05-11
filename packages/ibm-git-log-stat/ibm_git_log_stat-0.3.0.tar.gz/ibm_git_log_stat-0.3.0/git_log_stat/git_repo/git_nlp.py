from collections import defaultdict

from docx import Document
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
from docx.shared import Pt
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt
from reportlab.lib.pagesizes import A4
from reportlab.lib.units import inch
from reportlab.pdfgen import canvas
from transformers import pipeline

from git_log_stat.app_logs.logger_service import IBMLogger

log = IBMLogger("GitNlpService").get_logger()


def parse_commits(commit_lines):
    parsed = []
    for line in commit_lines:
        parts = [p.strip() for p in line.split('|')]
        if len(parts) == 4:
            parsed.append({'date': parts[0], 'author': parts[1], 'hash': parts[2], 'message': parts[3]})
    return parsed


summarizer = pipeline("summarization", model="facebook/bart-large-cnn")


def summarize_by_author(parsed_commits):
    author_msgs = defaultdict(list)
    for commit in parsed_commits:
        author_msgs[commit['author']].append(commit['message'])

    author_summaries = {}
    for author, messages in author_msgs.items():
        joined = " ".join(messages)
        if len(joined.split()) > 30:  # BART's min_length requirement
            summary = summarizer(joined, max_length=80, min_length=30, do_sample=False)[0]['summary_text']
        else:
            summary = "Too few commits for summarization. Messages: " + "; ".join(messages)
        author_summaries[author] = summary
    return author_summaries


def generate_natural_summary(commit_lines):
    # Combine all commit messages with authors
    combined_messages = []
    for line in commit_lines:
        parts = [p.strip() for p in line.split('|')]
        if len(parts) >= 3:
            date, author, *message_parts = parts
            message = " | ".join(message_parts)
            combined_messages.append(f"{author} - {message}")

    # Prepare input for the model
    text = ". ".join(combined_messages)

    # Optional: Truncate if too long
    if len(text.split()) > 1024:
        text = " ".join(text.split()[:1024])

    # Generate summary
    result = summarizer(text, max_length=100, min_length=30, do_sample=False)
    return result[0]['summary_text']


def generate_commit_summary_pdf(commit_logs: str, output_pdf_path="commit_summary.pdf"):
    # Clean and parse
    grouped_commits = defaultdict(list)
    combined_messages = []

    for line in commit_logs.split("\n"):
        parts = [p.strip() for p in line.split('|')]
        if len(parts) >= 4:
            date, author, commit_id, *message_parts = parts
            message = " | ".join(message_parts)
            grouped_commits[author].append((date, message))
            combined_messages.append(f"{author} - {message}")

    full_text = ". ".join(combined_messages)
    if len(full_text.split()) > 1024:
        full_text = " ".join(full_text.split()[:1024])

    summary = summarizer(full_text, max_length=120, min_length=40, do_sample=False)[0]['summary_text']

    # Step 2: PDF Generation
    c = canvas.Canvas(output_pdf_path, pagesize=A4)
    width, height = A4
    x_margin = 1 * inch
    y = height - 1 * inch

    def draw_multiline(text, start_y, font="Helvetica", size=12):
        c.setFont(font, size)
        for line in text.split('\n'):
            if start_y < inch:
                c.showPage()
                start_y = height - inch
                c.setFont(font, size)
            c.drawString(x_margin, start_y, line)
            start_y -= 14
        return start_y

    c.setFont("Helvetica-Bold", 16)
    c.drawString(x_margin, y, "Commit Summary Report")
    y -= 30

    c.setFont("Helvetica-Bold", 14)
    c.drawString(x_margin, y, "📝 Natural Language Summary:")
    y -= 20
    y = draw_multiline(summary, y)

    y -= 20
    c.setFont("Helvetica-Bold", 14)
    c.drawString(x_margin, y, "📜 Commits by Author:")
    y -= 20

    for author, commits in grouped_commits.items():
        c.setFont("Helvetica-Bold", 13)
        c.drawString(x_margin, y, f"👤 {author}")
        y -= 18

        for date, message in commits:
            is_pr = 'pull request' in message.lower() or 'merge' in message.lower()
            formatted = f"🔀 {date} | {message}" if is_pr else f"{date} | {message}"
            font_style = "Helvetica-Bold" if is_pr else "Helvetica"
            c.setFont(font_style, 11)
            y = draw_multiline(formatted, y)

        y -= 10  # Gap between authors

    c.save()
    log.info(f"✅ PDF saved to {output_pdf_path}")


def generate_commit_summary_docx(commit_logs: str, output_path="commit_summary.docx"):
    grouped_commits = defaultdict(list)
    combined_messages = []

    for line in commit_logs.split("\n"):
        parts = [p.strip() for p in line.split('|')]
        if len(parts) >= 4:
            date, author, commit_id, *message_parts = parts
            message = " | ".join(message_parts)
            grouped_commits[author].append((date, message))
            combined_messages.append(f"{author} - {message}")

    full_text = ". ".join(combined_messages)
    if len(full_text.split()) > 1024:
        full_text = " ".join(full_text.split()[:1024])

    summary = summarizer(full_text, max_length=120, min_length=40, do_sample=False)[0]['summary_text']

    # Step 2: Create Word document
    doc = Document()

    # Title
    title = doc.add_heading("Commit Summary Report", 0)
    title.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER

    doc.add_paragraph()  # Spacer

    # Natural Language Summary Section
    doc.add_heading("📝 Natural Language Summary", level=1)
    p_summary = doc.add_paragraph(summary)
    p_summary.style.font.size = Pt(11)

    doc.add_paragraph()  # Spacer

    # Commits grouped by author
    doc.add_heading("📜 Commits by Author", level=1)

    for author, commits in grouped_commits.items():
        doc.add_heading(f"👤 {author}", level=2)

        for date, message in commits:
            is_pr = "pull request" in message.lower() or "merge" in message.lower()
            formatted = f"🔀 {date} | {message}" if is_pr else f"{date} | {message}"
            p = doc.add_paragraph()
            run = p.add_run(formatted)
            run.font.size = Pt(10)
            if is_pr:
                run.bold = True

        doc.add_paragraph()  # Spacer between authors

    # Save the document
    doc.save(output_path)
    log.info(f"✅ DOCX saved to {output_path}")


def generate_commit_summary_pptx(commit_logs: str, output_file="commit_summary.pptx", skip_summary=False):
    # Initialize presentation
    prs = Presentation()

    # Title Slide
    slide_layout = prs.slide_layouts[0]  # Title slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Commit Summary Report"
    subtitle.text = "Detailed Report of Commits and Pull Requests"

    # Natural Language Summary Slide
    if not skip_summary:
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = "Natural Language Summary"

        # Load summarizer
        print("🔄 Summarizing commit logs...")

        # Combine messages for summarization
        combined_text = ". ".join([line.split("|", 2)[-1].strip() for line in commit_logs.split("\n")])
        combined_text = " ".join(combined_text.split()[:1024])  # to stay within token limits
        summary = summarizer(combined_text, max_length=255, min_length=40, do_sample=False)[0]['summary_text']

        # Add summary to the slide
        text_box = slide.shapes.placeholders[1]
        text_box.text = summary

    # Commits Grouped by Author Slide
    grouped_commits = defaultdict(list)
    for line in commit_logs.split("\n"):
        parts = [p.strip() for p in line.split('|')]
        if len(parts) >= 4:
            date, author, commit_id, *message_parts = parts
            message = " | ".join(message_parts)
            grouped_commits[author].append((date, commit_id, message))

    # Create a slide for each author and their commits
    for author, commits in grouped_commits.items():
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = f"Commits by {author}"

        # Add commit messages as bullet points
        text_box = slide.shapes.placeholders[1]
        text_frame = text_box.text_frame
        for date, commit_id, message in commits:
            p = text_frame.add_paragraph()
            p.text = f"{date} | {commit_id} | {message}"
            p.font.size = Pt(10)
            p.alignment = PP_ALIGN.LEFT

    # Save the PowerPoint file
    prs.save(output_file)
    print(f"✅ PowerPoint file saved to {output_file}")
