"""CLI for litdb.

The main command is litdb. There are subcommands for the actions.
"""

import os
import datetime
import json
import pathlib
import sys
import warnings

import bs4
import click
import dateparser
from docx import Document
from IPython import get_ipython
from IPython.display import display, HTML
from jinja2 import Template
from more_itertools import batched
import nbformat
from nbconvert import MarkdownExporter
import numpy as np

from pptx import Presentation
import requests
from rich import print as richprint
from rich.console import Console
from rich.markdown import Markdown
from sentence_transformers import SentenceTransformer
import tabulate
from tqdm import tqdm
import webbrowser


from .utils import get_config, init_litdb
from .db import get_db, add_source, add_work, add_author, update_filter, add_bibtex
from .openalex import get_data
from .pdf import add_pdf
from .bibtex import dump_bibtex
from .youtube import get_youtube_doc
from .audio import is_audio_url, get_audio_text, record
from .images import add_image, image_query, image_extensions

from .crawl import spider
from .research import deep_research
from .lsearch import llm_oa_search

with warnings.catch_warnings():
    warnings.filterwarnings("ignore", category=UserWarning, module="pydantic")
    from .gpt import gpt
    from .chat import chat

import logging
from transformers.utils import logging as tulogging

# Disable all Transformers logging
tulogging.set_verbosity_error()

logging.getLogger("pydantic").setLevel(logging.CRITICAL)

warnings.filterwarnings("ignore")
warnings.filterwarnings("ignore", category=UserWarning)

db = get_db()


@click.group()
def cli():
    """Group command for litdb."""
    pass


@cli.command()
def init():
    """Initialize a litdb directory in the current working directory."""
    init_litdb()
    db = get_db()

    with click.Context(about) as ctx:
        ctx.invoke(about)

    return db


#################
# Add functions #
#################


@cli.command()
@click.argument("sources", nargs=-1)
@click.option("--references", is_flag=True, help="Add references too.")
@click.option("--related", is_flag=True, help="Add related too.")
@click.option("--citing", is_flag=True, help="Add citing too.")
@click.option("--all", is_flag=True, help="Add references, related and citing.")
@click.option("-t", "--tag", "tags", multiple=True)
def add(
    sources,
    references=False,
    citing=False,
    related=False,
    all=False,
    verbose=False,
    tags=None,
):
    """Add WIDS to the db.

    REFERENCES, RELATED, CITING are flags to also add those for DOI sources. ALL
    is shorthand for all of those.

    SOURCES can be one or more of a doi or orcid, a pdf path, a url, bibtex
    file, or other kind of file assumed to be text.

    TAGS is a list of tags to add to the source.

    These are one time additions.

    """

    for source in tqdm(sources):
        # a work
        if source.startswith("10.") or "doi.org" in source:
            if source.startswith("10."):
                source = f"https://doi.org/{source}"

            if all:
                references, citing, related = True, True, True

            add_work(source, references, citing, related)

        # works from an author
        elif "orcid" in source or source.lower().startswith("https://openalex.org/a"):
            add_author(source)

        # a bibtex file
        elif source.endswith(".bib"):
            add_bibtex(source)

        # pdf
        elif source.endswith(".pdf"):
            source = os.path.abspath(source)
            add_pdf(source)

        # docx
        elif source.endswith(".docx"):
            source = os.path.abspath(source)
            doc = Document(source)
            add_source(source, "\n".join([para.text for para in doc.paragraphs]))

        # pptx
        elif source.endswith(".pptx"):
            source = os.path.abspath(source)
            prs = Presentation(source)
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            add_source(source, "\n".join(text))

        # YouTube
        elif source.startswith("https") and "youtube" in source:
            text, citation = get_youtube_doc(source)
            add_source(source, text, {"citation": citation})

        # audio sources
        elif (source.startswith("http") and is_audio_url(source)) or source.endswith(
            ".mp3"
        ):
            add_source(source, get_audio_text(source))

        # local html
        elif not source.startswith("http") and source.endswith(".html"):
            source = os.path.abspath(source)
            with open(source) as f:
                text = f.read()
            soup = bs4.BeautifulSoup(text, features="lxml")
            add_source(source, soup.get_text())

        # a url
        elif source.startswith("http"):
            soup = bs4.BeautifulSoup(requests.get(source).text)
            add_source(source, soup.get_text())

        # ipynb
        elif source.endswith(".ipynb"):
            source = os.path.abspath(source)
            with open(source) as f:
                notebook = nbformat.read(f, as_version=4)

            # Create a Markdown exporter
            markdown_exporter = MarkdownExporter()

            # Convert the notebook to Markdown
            (body, resources) = markdown_exporter.from_notebook_node(notebook)

            add_source(source, body)

        # There are a lot of image extensions. I put this near the end so the
        # specific extensions are matched first.
        elif os.path.splitext(source)[1].lower() in image_extensions:
            add_image(source)

        # assume it is text
        else:
            source = os.path.abspath(source)
            with open(source) as f:
                text = f.read()
            add_source(source, text)

    if tags:
        with click.Context(add_tag) as ctx:
            ctx.invoke(add_tag, sources=sources, tags=tags)


@cli.command()
@click.argument("sources", nargs=-1)
def remove(sources):
    """Remove sources from litdb."""
    for source in sources:
        db.execute("delete from sources where source = ?", (source,))
        db.commit()


@cli.command()
@click.argument("query", nargs=-1)
@click.option("--references", is_flag=True)
@click.option("--related", is_flag=True)
@click.option("--citing", is_flag=True)
def crossref(query, references, related, citing):
    """Add entries to litdb from a crossref query."""
    query = " ".join(query)
    resp = requests.get("https://api.crossref.org/works", params={"query": query})

    if resp.status_code == 200:
        data = resp.json()
        for i, item in enumerate(data["message"]["items"]):
            authors = ", ".join(
                [f'{au["given"]} {au["family"]}' for au in item.get("author", [])]
            )
            source = " ".join(item.get("container-title", ["no source"]))
            published = item.get("published", {}) or {}
            year = published.get("date-parts", [["no year"]])[0][0]
            title = item.get("title", ["no title"])
            richprint(
                f'{i}. {" ".join(title)}, {authors}, {source} ({year}), https://doi.org/{item["DOI"]}.'
            )

        toadd = input("Enter space separated numbers to add, or return to quit. ")

        if toadd:
            toadd = [int(x) for x in toadd.split(" ")]
            dois = [
                "https://doi.org/" + data["message"]["items"][i]["DOI"] for i in toadd
            ]

            with click.Context(add) as ctx:
                ctx.invoke(
                    add,
                    sources=dois,
                    related=related,
                    references=references,
                    citing=citing,
                )


@cli.command()
@click.argument("sources", nargs=-1)
def index(sources):
    """Index the directories in SOURCES.

    SOURCES is a list of directories.
    """
    for directory in sources:
        directory = pathlib.Path(directory).resolve()
        for fname in directory.rglob("*"):
            # for f in files:
            if fname.suffix in [
                ".pdf",
                ".docx",
                ".pptx",
                ".org",
                ".md",
                ".html",
                ".bib",
                ".ipynb",
            ]:
                fname = str(fname)

                # skip files we already have
                if db.execute(
                    """select source from sources where source = ?""", (fname,)
                ).fetchone():
                    continue

                with click.Context(add) as ctx:
                    print(fname)
                    ctx.invoke(add, sources=[fname])
                    print(f"Adding {fname}")

                    richprint(f"Added {fname}")

        last_updated = datetime.date.today().strftime("%Y-%m-%d")

        directory = str(directory)  # we need strings for the db
        if db.execute(
            """select path from directories where path = ?""", (directory,)
        ).fetchone():
            print(f"Updating {directory}")
            db.execute(
                """update directories set last_updated = ?
            where path = ?""",
                (last_updated, directory),
            )
        else:
            print(f"Inserting {directory}: {last_updated}")
            db.execute(
                """insert into directories(path, last_updated)
            values (?, ?)""",
                (directory, last_updated),
            )

        db.commit()


@cli.command()
def reindex():
    """Reindex saved directories."""
    for (directory,) in db.execute("""select path from directories""").fetchall():
        print(f"Reindexing {directory}")
        index([directory])


###########
# Tagging #
###########
@cli.command()
@click.argument("sources", nargs=-1)
@click.option("-t", "--tag", "tags", multiple=True)
def add_tag(sources, tags):
    """Add tags to sources.

    It is a little annoying to add multiple tags. It looks like this.
    litdb add-tag source -t tag1 -t tag2
    """
    for source in sources:
        # Get source id
        (source_id,) = db.execute(
            "select rowid from sources where source = ?", (source,)
        ).fetchone()

        for tag in tags:
            # get tag id
            tag_id = db.execute(
                "select rowid from tags where tag = ?", (tag,)
            ).fetchone()

            if not tag_id:
                c = db.execute("insert into tags(tag) values (?)", (tag,))
                tag_id = c.lastrowid
                db.commit()
            else:
                # we get a tuple in the first query
                (tag_id,) = tag_id

            # Now add a tag
            db.execute(
                "insert into source_tag(source_id, tag_id) values (?, ?)",
                (source_id, tag_id),
            )
            db.commit()

            print(f"Tagged {source} with {tag}")


@cli.command()
@click.argument("sources", nargs=-1)
@click.option("-t", "--tag", "tags", multiple=True)
def rm_tag(sources, tags):
    """Remove tags from sources.

    It is a little annoying to remove multiple tags. It looks like this.
    litdb rm-tag source -t tag1 -t tag2
    """
    for source in sources:
        # Get source id
        (source_id,) = db.execute(
            "select rowid from sources where source = ?", (source,)
        ).fetchone()

        for tag in tags:
            # get tag id. Assume it exists?
            (tag_id,) = db.execute(
                "select rowid from tags where tag = ?", (tag,)
            ).fetchone()

            c = db.execute(
                """delete from source_tag
            where source_id = ? and tag_id = ?""",
                (source_id, tag_id),
            )

            db.commit()
            print(f"Deleted {c.rowcount} rows ({tag} from {source}")


@cli.command()
@click.argument("tags", nargs=-1)
def delete_tag(tags):
    """Delete each tag.

    This should also delete tags from sources by cascade.
    """
    for tag in tags:
        c = db.execute("delete from tags where tag = ?", (tag,))
        print(f"Deleted {c.rowcount} rows ({tag})")
    db.commit()


@cli.command()
@click.argument("tags", nargs=-1)
@click.option("-f", "--fmt", default='{{ source }}\n{{ extra["citation"] }}')
def show_tag(tags, fmt):
    """Show entries with tags.

    FMT is a jinja template for the output. You have variables of source, text
    and extra.

    I don't have good logic here, we just show all entries. I could probably get
    some basic and logic with sets, but mostly I assume for now you only want
    one tag, so this works. TODO: add something like boolean logic?

    """
    template = Template(fmt)
    for tag in tags:
        for row in db.execute(
            """select
        sources.source, sources.text, sources.extra
        from sources
        inner join source_tag on source_tag.source_id = sources.rowid
        inner join tags on source_tag.tag_id = tags.rowid
        where tags.tag = ?""",
            (tag,),
        ).fetchall():
            source, text, extra = row
            extra = json.loads(extra)
            richprint(template.render(**locals()))


@cli.command()
def list_tags():
    """Print defined tags."""
    print("The following tags are defined.")
    for (tag,) in db.execute("select tag from tags").fetchall():
        print(tag)


##########
# Review #
##########


@cli.command()
@click.option("-s", "--since", default="1 week ago")
@click.option("-f", "--fmt", default=None)
def review(since, fmt):
    """Review new entries added SINCE.

    SINCE should be something dateparser can handle.
    FMT is a jinja template for the output. Defaults to an org-mode template.
    """

    since = dateparser.parse(since).strftime("%Y-%m-%d")
    c = db.execute(
        """select source, text, extra from sources
    where date(date_added) > ?""",
        (since,),
    ).fetchall()

    template = Template(
        fmt
        or """* {{ extra['display_name'] | replace("\n", " ") }}
:PROPERTIES:
:SOURCE: {{ source }}
:OPENALEX: {{ extra.get('id') }}
:YEAR: {{ extra.get('publication_year') }}
:REFERENCE_COUNT: {{ extra.get('referenced_works_count', 0) }}
:CITED_BY_COUNT: {{ extra.get('cited_by_count', 0) }}
:END:

{{ text }} litdb:{{ source }}
        """
    )

    for source, text, extra in c:
        extra = json.loads(extra) or {}
        print(template.render(**locals()))


#############
# Searching #
#############


@cli.command()
def screenshot():
    """Do vector search from text in a screenshot.

    Use OCR to get text from an image on the clipboard (probably from a
    screenshot) and do a vector search on the text.
    """
    from PIL import ImageGrab
    import pytesseract

    os.environ["TOKENIZERS_PARALLELISM"] = "false"

    # Grab the image from the clipboard
    img = ImageGrab.grabclipboard()

    if img:
        text = pytesseract.image_to_string(img)
        print(f'Searching for "{text}"')
        vsearch(text)
    else:
        print("No image found in clipboard.")


@cli.command()
@click.option("-p", "--playback", is_flag=True, help="Play audio back")
def audio(playback=False):
    """Record audio, convert it to text, and do a vector search on the text.

    The idea is nice, but the quality of transcription for scientific words is
    not that great. A better transcription library might make this more useful.
    """

    while True:
        afile = record()
        text = get_audio_text(afile)
        print("\n" + text + "\n")

        if playback:
            import playsound

            playsound.playsound(afile, block=True)

        response = input("Is that what you want to search? ([y]/n/q): ")
        if response.lower().startswith("q"):
            return
        elif response.lower().startswith("n"):
            # record a new audio
            continue
        else:
            # move on to searching
            break

    vsearch([text])


@cli.command()
@click.argument("query", nargs=-1)
@click.option("-n", default=3)
@click.option("-e", "--emacs", is_flag=True, default=False)
@click.option("-i", "--iterative", is_flag=True, default=False)
@click.option("-m", "--max-steps", default=None)
@click.option(
    "-f",
    "--fmt",
    default=(" {{ i }}. ({{ score|round(3) }})" " {{ source }}\n{{ text[:200] }}\n\n"),
)
@click.option("-x", "--cross-encode", is_flag=True, default=False)
def vsearch(query, n, emacs, fmt, cross_encode, iterative, max_steps):
    """Do a vector search on QUERY.

    N is an integer for number of results to return
    EMACS is a flag for changing the output format for emacs

    FORMAT is a jinja template for the output. The variables you have access to
    are i, source, text, extra, similarity.

    CROSS_ENCODE is a boolean that resorts the results with a cross-encoder.

    ITERATIVE is a boolean that expands the search by references, citations and
    related articles from the top matches until you tell it to stop or reach
    MAX_STEPS.

    MAX_STEPS is the maximum number of iterations to take. If you set this, you
    will not be prompted each time, it will just run those steps until nothing
    better is found, or you reach the number.

    """
    config = get_config()
    query = " ".join(query)
    model = SentenceTransformer(config["embedding"]["model"])
    emb = model.encode([query]).astype(np.float32).tobytes()

    if iterative:
        best = None

        steps = 0

        while True:
            results = db.execute(
                """select
            sources.source, sources.text,
            sources.extra, vector_distance_cos(?, embedding) as d
            from vector_top_k('embedding_idx', ?, ?)
            join sources on sources.rowid = id
            order by d""",
                (emb, emb, n),
            ).fetchall()

            for source, text, extra, d in results:
                richprint(f"{d:1.3f}: {source}")

            steps += 1
            if steps == max_steps:
                break

            current = [x[0] for x in results]  # sources

            # This means no change
            if current == best:
                print("Nothing new was found")
                break

            if not max_steps and input(
                "Search for better matches? ([y]/n)"
            ).lower().startswith("n"):
                break

            # something changed. add references and loop
            best = current
            for source in current:
                add_work(source, True, True, True)

    else:
        c = db.execute(
            """select sources.source, sources.text,
        sources.extra, vector_distance_cos(?, embedding)
        from vector_top_k('embedding_idx', ?, ?)
        join sources on sources.rowid = id""",
            (emb, emb, n),
        )

        results = c.fetchall()

    if cross_encode:
        import torch
        from sentence_transformers.cross_encoder import CrossEncoder

        # I don't know why I have to set the activation function here, but the
        # score is not 0..1 otherwise
        ce = CrossEncoder(
            config["embedding"]["cross-encoder"],
            default_activation_function=torch.nn.Sigmoid(),
        )
        scores = ce.predict([[query, text] for _, text, _, _ in results])
        # resort based on the scores
        results = [results[i] for i in np.argsort(scores)]

    if emacs:
        tmpl = (
            "( {% for source, text, extra, score in results %}"
            '("({{ score|round(3) }}) {{ text }}" . "{{ source }}") '
            "{% endfor %})"
        )
        template = Template(tmpl)
        print(template.render(**locals()))
    else:
        for i, row in enumerate(results, 1):
            source, text, extra, score = row
            template = Template(fmt)
            richprint(template.render(**locals()))

    return results


@cli.command()
@click.argument("query", nargs=-1)
@click.option("-n", default=3)
@click.option(
    "-f", "--fmt", default="{{ source }} ({{ score | round(3) }})\n{{ snippet }}"
)
def fulltext(query, n, fmt):
    """Perform a fulltext search on litdb."""
    query = " ".join(query)

    results = db.execute(
        """select
    sources.source, sources.text, snippet(fulltext, 1, '', '', '', 16), sources.extra, bm25(fulltext)
    from fulltext
    inner join sources on fulltext.source = sources.source
    where fulltext.text match ? order by rank limit ?""",
        (query, n),
    ).fetchall()

    for source, text, snippet, extra, score in results:
        richprint(Template(fmt).render(**locals()))

    return results


gpt = cli.command(gpt)


@click.command(help=chat.__doc__)
@click.option("--model", default=None, help="The LiteLLM model to use.")
@click.option("--debug", is_flag=True, default=False)
def chat_command(model, debug):
    chat(model, debug)


cli.add_command(chat_command, name="chat")


@cli.command(help=spider.__doc__)
@click.argument("root")
def crawl(root):
    """Crawl a website at ROOT url."""
    spider(root)


@cli.command()
@click.argument("query", nargs=-1)
@click.option("-q", default=5, help="The number of queries to generate")
@click.option("-n", default=25, help="The number of results to get for each query")
@click.option("-k", default=5, help="The number of results to return")
def lsearch(query, q, n, k):
    """LLM enhanced search of OpenAlex.

    QUERY: string, an natural language query
    Q: int, number of keyword searches to generate
    N: int, number of results to retrieve for each keyword query
    K: int, number of results to finally return

    Internally, it generates Q keyword queries based on the original QUERY using
    an LLM. For each keyword query several searches are run sorted on citations
    and publication year both ascending and descending, and a random sample is
    searched. Then these are combined and sorted by vector similarity to the
    query. Finally the top k results are printed.

    This does not add anything to the litdb database.

    """
    for s, result in llm_oa_search(query, q, n, k):
        richprint(
            f"{s[0]:1.2f}: {result['title']} ({result['publication_year']}), {result['id']}\n\n"
        )

from .extract import extract_tables, extract_schema

@cli.command()
@click.argument('pdf', type=click.Path(exists=True, dir_okay=False, readable=True), required=True)
@click.option('-t', '--tables', type=int, multiple=True,
              help="Table numbers to extract (1-based index).")
@click.option('-f', '--fmt', default='csv')
def extract(pdf, tables, fmt):
    """Extract tables from a pdf.

    PDF: string, path to file
    TABLES: list of int, the table numbers to extract, starting at 1
    FMT: string, output format.
    """

    for df in extract_tables(pdf, tables):
        match fmt:
            case "csv":
                print(df.to_csv(index=False))
            case "json":
                print(df.to_json(index=False))
            case "md":
                print(df.to_markdown(index=False))
            case _:
                print(df)
        print()


@cli.command()
@click.argument('source', type=str, required=True)
@click.argument('schema', type=str, required=True)
def schema(source, schema):
    """Extract structured schema from a SOURCE.

    SOURCE: string, url or path to file
    SCHEMA: string, the scheme to extract. Comma separated name type.
    """
    print(extract_schema(source, schema))


@cli.command()
@click.argument("query", nargs=-1)
@click.option(
    "--report-type", default="research_report", help="The type of report to generate."
)
@click.option("--doc-path", default=None, help="Path to local documents")
@click.option("-o", "--output", default=None, help="output file")
@click.option("-v", "--verbose", is_flag=True, default=False)
def research(query, report_type, doc_path, output, verbose):
    """Run a deep research query.

    QUERY: the topic to do research on

    REPORT_TYPE: one of the supported types in gpt_researcher
    DOC_PATH: a directory path for local files
    OUTPUT: a filename to write output to, defaults to printing to stdout
    VERBOSE: if truthy the output is more verbose.

    Based on gpt_researcher. You need to have some configuration setup in advance.
    API keys for the LLM in environment variables.
    """
    query = " ".join(query)  # if you don't quote the query it is a list of words.

    if doc_path:
        os.environ["DOC_PATH"] = doc_path

    report, result, context, costs, images, sources = deep_research(
        query, report_type, verbose
    )

    s = f"""{report}

# Research costs
${costs}

# Result
{result}

# Context
{context}
"""
    if output:
        base, ext = os.path.splitext(output)
        # I found pypandoc was not good at pdf. lots of bad latex commands that
        # make the pdf build fail.
        # pdfkit relies on wkhtmltopdf which appears discontinued
        # weasyprint and m2pdf has some gobject dependency
        # These are adapted from gpt_researcher / multi_agents

        if ext == ".pdf":
            from md2pdf.core import md2pdf

            md2pdf(output, md_content=s)

        elif ext == ".docx":
            import mistune
            from htmldocx import HtmlToDocx
            from docx import Document

            html = mistune.html(s)
            doc = Document()
            HtmlToDocx().add_html_to_document(html, doc)
            doc.save(output)

        elif ext == ".html":
            import mistune

            html = mistune.html(s)
            with open(output, "w") as f:
                f.write(html)

        elif ext == ".org":
            import pypandoc

            with open(output, "w") as f:
                org = pypandoc.convert_text(s, to="org", format="md")
                f.write(org)

        elif ext == ".md":
            with open(output, "w") as f:
                f.write(s)
        else:
            print(f"I do not know how to make {output}.")

        import webbrowser

        if os.path.exists(output):
            print(f"Opening {output}")
            webbrowser.open(f"file://{os.path.abspath(output)}")

    else:
        console = Console(color_system="truecolor")

        with console.pager():
            console.print(Markdown(s))


@cli.command()
@click.argument("query", nargs=-1)
@click.option("-c", "--clipboard", is_flag=True, default=False)
@click.option("-n", default=1, help="Number of hits to retrieve")
def image_search(query, clipboard, n):
    image_query(" ".join(query), clipboard, n)


@cli.command()
@click.argument("source")
@click.option("-f", "--fmt", default=None)
@click.option("-e", "--emacs", is_flag=True, default=False)
@click.option("-n", default=3)
def similar(source, n, emacs, fmt):
    """Find N sources similar to SOURCE by vector similarity.

    if EMACS is truthy, the output is lisp for Emacs to read.
    FMT is a jinja template with access to source, text, and extra
    """
    (emb,) = db.execute(
        """select embedding from sources where source = ?""", (source,)
    ).fetchone()

    allrows = db.execute(
        """select sources.source, sources.text, sources.extra
    from vector_top_k('embedding_idx', ?, ?)
    join sources on sources.rowid = id""",
        # we do n + 1 because the first entry is always the source
        (emb, n + 1),
    ).fetchall()

    rows = [(source, text, json.loads(extra)) for source, text, extra in allrows[1:]]

    if emacs:
        template = Template(
            "({% for source, text, extra in rows %}"
            ' ("{{ extra.get("citation") or text }}" . "{{ source }}")'
            " {% endfor %})"
        )
        print(template.render(**locals()))
    else:
        template = Template(fmt or "{{ i }}. {{ source }}\n {{text}}\n\n")
        # print starting at index 1, the first item is always the source.
        for i, row in enumerate(rows, 1):
            source, text, extra = row
            richprint(template.render(**locals()))


@cli.command()
@click.argument("vector_query")
@click.argument("text_query")
@click.option("-n", default=5)
@click.option(
    "-f", "--fmt", default="{{ source }} ({{ score | round(3) }})\n{{ text }}\n\n"
)
def hybrid_search(vector_query, text_query, n, fmt):
    """Perform a hybrid vector and full text search.

    VECTOR_QUERY: The query to do vector search on
    TEXT_QUERY: The query for full text search
    N is an integer number of documents to return.
    FMT is a jinja template.
    """
    # Get vector results and score
    with click.Context(vsearch) as ctx:
        # source, text, extra, score
        vresults = ctx.invoke(vsearch, query=vector_query.split(" "), n=n, fmt="")
        vscores = [(result[0], result[3]) for result in vresults]

    # Get text results and score
    with click.Context(fulltext) as ctx:
        tresults = ctx.invoke(fulltext, query=text_query.split(" "), n=n, fmt="")
        # I think sqlite makes scores negative to sort them the way they want. I
        # reverse this here.
        tscores = [(result[0], -result[-1]) for result in tresults]

    # Normalize scores
    minv, maxv = min([x[1] for x in vscores]), max([x[1] for x in vscores])
    mint, maxt = min([x[1] for x in tscores]), max([x[1] for x in tscores])

    vscores = {oaid: (score - minv) / (maxv - minv) for oaid, score in vscores}
    tscores = {oaid: (score - minv) / (maxv - minv) for oaid, score in tscores}

    combined_scores = {}
    for oaid in set(vscores.keys()).union(tscores.keys()):
        vscore = vscores.get(oaid, 0)
        tscore = tscores.get(oaid, 0)
        cscore = 1 / (1 + vscore) + 1 / (1 + tscore)
        combined_scores[oaid] = cscore

    sorted_results = sorted(combined_scores.items(), key=lambda x: x[1], reverse=True)
    results = []
    for oaid, score in sorted_results:
        c = db.execute(
            "select source, text, extra from sources where source = ?", (oaid,)
        )
        row = c.fetchone()
        results += [[*row, score]]

    for row in results:
        source, text, extra, score = row
        richprint(Template(fmt).render(**locals()))

    return results


@cli.command()
@click.argument("query", nargs=-1)
@click.option("-g", "--google", is_flag=True)
@click.option("-gs", "--google-scholar", is_flag=True)
@click.option("-pm", "--pubmed", is_flag=True)
@click.option("-ar", "--arxiv", is_flag=True)
@click.option("-cr", "--chemrxiv", is_flag=True)
@click.option("-br", "--biorxiv", is_flag=True)
@click.option("-a", "--all", is_flag=True)
def web(query, google, google_scholar, pubmed, arxiv, chemrxiv, biorxiv, all):
    """Open a web search for QUERY.

    We always do an OpenAlex search.

    If these are true, we also search here
    GOOGLE
    GOOGLE_SCHOLAR
    PUBMED
    ARXIV
    CHEMRXIV
    BIORXIV
    """
    query = " ".join(query)

    if all:
        google, google_scholar = (
            True,
            True,
        )
        pubmed, arxiv, chemrxiv, biorxiv = True, True, True, True

    # This is to avoid some huggingface/tokenizer warning. I don't know why we
    # need to do it, but this code forks the process, and triggers that warning.
    os.environ["TOKENIZERS_PARALLELISM"] = "false"

    oa = f"https://openalex.org/works?filter=default.search:{query}"
    webbrowser.open(oa)

    if google:
        url = f"https://www.google.com/search?q={query}"
        webbrowser.open(url)

    if google_scholar:
        url = f"https://scholar.google.com/scholar?q={query}"
        webbrowser.open(url)

    if pubmed:
        url = f"https://pubmed-ncbi-nlm-nih-gov.cmu.idm.oclc.org/?term={query}"
        webbrowser.open(url)

    if arxiv:
        url = f"https://arxiv.org/search/?query={query}"
        webbrowser.open(url)

    if chemrxiv:
        url = f"https://chemrxiv.org/engage/chemrxiv/search-dashboard?text={query}"
        webbrowser.open(url)

    if biorxiv:
        url = f"https://www.biorxiv.org/search/{query}"
        webbrowser.open(url)


###########
# Filters #
###########


@cli.command()
@click.argument("_filter")
@click.option("-d", "--description")
def add_filter(_filter, description=None):
    """Add an OpenAlex FILTER.

    This does not run the filter right away. You need
    to manually update the filters later.
    """
    db.execute(
        "insert into queries(filter, description) values (?, ?)", (_filter, description)
    )
    db.commit()


@cli.command()
@click.argument("_filter")
def rm_filter(_filter):
    """Remove an OpenAlex FILTER."""
    db.execute("delete from queries where filter = ?", (_filter,))
    db.commit()


update_filter_fmt = """** {{ extra['display_name'] | replace("\n", "") | replace("\r", "") }}
:PROPERTIES:
:SOURCE: {{ source }}
:REFERENCE_COUNT: {{ extra.get('referenced_works_count', 0) }}
:CITED_BY_COUNT: {{ extra.get('cited_by_count', 0) }}
:END:

litdb:{{ source }}

{{ text }}

"""


@cli.command()
@click.option("-f", "--fmt", default=update_filter_fmt)
@click.option("-s", "--silent", is_flag=True, default=False)
def update_filters(fmt, silent):
    """Update litdb using a filter with works from a created date."""
    filters = db.execute("""select filter, description, last_updated from queries""")
    for f, description, last_updated in filters.fetchall():
        results = update_filter(f, last_updated, silent)
        if results:
            richprint(f"* {description}")
        for result in results:
            source, text, extra = result
            richprint(Template(fmt).render(**locals()))


list_filter_fmt = (
    '{{ "{:3d}".format(rowid) }}.'
    ' {{ "{:30s}".format(description'
    ' or "None") }} {{ f }}'
    " ({{ last_updated }})"
)


@cli.command()
@click.option("-f", "--fmt", default=list_filter_fmt)
def list_filters(fmt):
    """List the filters.

    FMT is a jinja template with access to the variables rowid, f, description
    and last_updated. f is the filter string.

    You can dump the filters to stdout like this.

    > litdb list-filters -f 'litdb add-filter {{ f }}'

    You could use that to send a list of your filters to someone, or to recreate
    a db somewhere else.
    """
    filters = db.execute(
        """select rowid, filter, description, last_updated
    from queries"""
    )
    for rowid, f, description, last_updated in filters.fetchall():
        richprint(Template(fmt).render(**locals()))


######################
# OpenAlex searching #
######################


@cli.command()
@click.argument("query", nargs=-1)
@click.option("-f", "--filter", "_filter", is_flag=True, default=False)
@click.option("-e", "--endpoint", default="works")
@click.option("--sample", default=-1)
@click.option("--per-page", default=5)
def openalex(query, _filter, endpoint, sample, per_page):
    """Run an openalex query on FILTER.

    ENDPOINT should be one of works, authors, or another entity.
    SAMPLE: int, return this many random samples
    PER_PAGE: limits the number of results retrieved

    This does not add anything to your database. It is to help you find starting
    points.

    To search text:
    litdb openalex "circular polymer"

    To find a journal id with a specific filter
    litdb openalex -e sources -f "display_name.search:Digital Discovery"

    """
    config = get_config()
    url = f"https://api.openalex.org/{endpoint}"

    if isinstance(query, tuple):
        query = " ".join(query)
    if not _filter:
        query = f"default.search:{query}"

    params = {
        "email": config["openalex"]["email"],
        "filter": query,
        "per_page": per_page,
    }

    if api_key := config["openalex"].get("api_key"):
        params.update(api_key=api_key)

    if sample > 0:
        params.update(sample=sample, per_page=sample)

    resp = requests.get(url, params)
    if resp.status_code != 200:
        print(resp.url)
        print(resp.text)
        return

    data = resp.json()
    for result in data["results"]:
        s = f'{result["title"]}, ({result["publication_year"]}) {result["id"]}\n'
        # Note sometimes there is an exception from bad markup in strings
        richprint(s)


########################################
# Convenience functions to add filters #
########################################


@cli.command()
@click.argument("name", nargs=-1)
def author_search(name):
    """Search OpenAlex for name.

    Uses the autocomplete endpoint to find an author's orcid.
    """
    auname = " ".join(name)

    url = "https://api.openalex.org/autocomplete/authors"

    from .openalex import get_data

    data = get_data(url, params={"q": auname})

    for result in data["results"]:
        richprint(
            f'- {result["display_name"]}\n  {result["hint"]} '
            f'{result["external_id"]}\n\n'
        )


@cli.command()
@click.argument("orcids", nargs=-1)
@click.option("-r", "--remove", is_flag=True, help="remove")
def follow(orcids, remove=False):
    """Add a filter to follow orcid."""
    for orcid in orcids:
        if not orcid.startswith("http"):
            orcid = f"https://orcid.org/{orcid}"

        # Seems like we should get their articles first.
        add_author(orcid)

        f = f"author.orcid:{orcid}"

        if remove:
            c = db.execute("""delete from queries where  filter = ?""", (f,))
            db.commit()
            richprint(f"{c.rowcount} rows removed")
            return

        url = f"https://api.openalex.org/authors/{orcid}"
        data = get_data(url)
        name = data["display_name"]

        today = datetime.date.today().strftime("%Y-%m-%d")
        db.execute(
            """insert or ignore into
        queries(filter, description, last_updated)
        values (?, ?, ?)""",
            (f, name, today),
        )

        richprint(f"Following {name}: {orcid}")
        db.commit()


@cli.command()
@click.argument("query", nargs=-1)
@click.option("-r", "--remove", is_flag=True, help="remove")
def watch(query, remove=False):
    """Create a watch on query.

    QUERY: string, a filter for openalex.
    REMOVE: a flag to remove the query.
    """
    # First, we should make sure the filter is valid
    query = " ".join(query)

    if remove:
        c = db.execute("""delete from queries where filter = ?""", (query,))
        db.commit()
        richprint(f"{c.rowcount} rows removed")
        return

    url = "https://api.openalex.org/works"

    data = get_data(url, params={"filter": query})
    if len(data["results"]) == 0:
        richprint(f"Sorry, {query} does not seem valid.")

    if remove:
        c = db.execute("""delete from queries where filter = ?""", (query,))
        richprint(f"Deleted {c.rowcount} rows")
        db.commit()
    else:
        c = db.execute(
            """insert or ignore into queries(filter, description)
        values (?, ?)""",
            (query,),
        )
        richprint(f"Added {c.rowcount} rows")
        db.commit()
        richprint(f"Watching {query}")


@cli.command()
@click.argument("doi")
@click.option("-r", "--remove", is_flag=True, help="remove")
def citing(doi, remove=False):
    """Create a watch for articles that cite doi.

    REMOVE is a flag to remove the doi.
    """
    url = "https://api.openalex.org/works"

    # We need an OpenAlex id
    f = f"doi:{doi}"

    data = get_data(url, params={"filter": f})
    if len(data["results"]) == 0:
        richprint(f"Sorry, {doi} does not seem valid.")

    wid = data["results"][0]["id"]

    if remove:
        c = db.execute("""delete from queries where filter = ?""", (f"cites:{wid}",))
        db.commit()
        richprint(f"Deleted {c.rowcount} rows")
    else:
        c = db.execute(
            """insert or ignore into queries(filter, description)
        values (?, ?)""",
            (f"cites:{wid}", f"Citing papers for {doi}"),
        )

        db.commit()
        richprint(f"Added {c.rowcount} rows")


@cli.command()
@click.argument("doi")
@click.option("-r", "--remove", is_flag=True, help="remove")
def related(doi, remove=False):
    """Create a watch for articles that are related to doi.

    REMOVE is a flag to remove the doi from queries.
    """
    url = "https://api.openalex.org/works"

    # We need an OpenAlex id
    f = f"doi:{doi}"

    data = get_data(url, params={"filter": f})
    if len(data["results"]) == 0:
        richprint(f"Sorry, {doi} does not seem valid.")

    wid = data["results"][0]["id"]

    if remove:
        c = db.execute(
            """delete from queries where filter = ?""", (f"related_to:{wid}",)
        )
        db.commit()
        richprint(f"Deleted {c.rowcount} rows")
    else:
        c = db.execute(
            """insert or ignore into queries(filter, description)
        values (?, ?)""",
            (f"related_to:{wid}", f"Related papers for {doi}"),
        )

        db.commit()
        richprint(f"Added {c.rowcount} rows")


#############
# Utilities #
#############


@cli.command()
@click.argument("sources", nargs=-1)
def bibtex(sources):
    """Generate bibtex entries for sources."""
    if not sources:
        sources = sys.stdin.read().strip().split()

    for source in sources:
        work = db.execute(
            """select extra from sources where source = ?""", (source,)
        ).fetchone()
        if work:
            richprint(dump_bibtex(json.loads(work[0])))
        else:
            print(f"No entry found for {source}")


@cli.command()
@click.argument("sources", nargs=-1)
def citation(sources):
    """Generate citation strings for sources."""
    if not sources:
        sources = sys.stdin.read().strip().split()

    for i, source in enumerate(sources):
        (_citation,) = db.execute(
            """select json_extract(extra, '$.citation')
        from sources where source = ?""",
            (source,),
        ).fetchone()
        richprint(f"{i + 1:2d}. {_citation}")


@cli.command()
@click.argument("doi")
def unpaywall(doi):
    """Use unpaywall to find PDFs for doi."""
    config = get_config()
    url = f"https://api.unpaywall.org/v2/{doi}"
    params = {"email": config["openalex"]["email"]}

    resp = requests.get(url, params)
    if resp.status_code == 200:
        data = resp.json()
        richprint(f'{data["title"]}, {data.get("journal_name") or ""}')
        richprint(f'Is open access: {data.get("is_oa", False)}')

        for loc in data.get("oa_locations", []):
            richprint(loc.get("url_for_pdf") or loc.get("url_for_landing_page"))
    else:
        richprint(f"{doi} not found in unpaywall")


@cli.command()
def about():
    """Summary statistics of your db."""
    config = get_config()
    dbf = os.path.join(config["root"], "litdb.libsql")
    cf = os.path.join(config["root"], "litdb.toml")

    richprint(f"Your database is located at {dbf}")
    richprint(f"The configuration is at {cf}")
    kb = 1024
    mb = 1024 * kb
    gb = 1024 * mb
    richprint(f"Database size: {os.path.getsize(dbf) / gb:1.2f} GB")
    db = get_db()
    (nsources,) = db.execute("select count(source) from sources").fetchone()
    richprint(f"You have {nsources} sources")


@cli.command()
@click.argument("sql")
def sql(sql):
    """Run the SQL command on the db."""
    for row in db.execute(sql).fetchall():
        richprint(row)


@cli.command()
@click.argument("sources", nargs=-1)
@click.option("-f", "--fmt", default="{{ source }}\n{{ text }}")
def show(sources, fmt):
    """Show the source.

    FMT is a jinja template with access to source, text, extra for each arg.
    """
    for src in sources:
        result = db.execute(
            """select source, text, extra from
        sources where source = ?""",
            (src,),
        ).fetchone()

        if result:
            source, text, extra = result
            extra = json.loads(extra)
            richprint(Template(fmt).render(**locals()))
        else:
            print(f"Nothing found for {src}")


@cli.command(name="open")
@click.argument("source")
def visit(source):
    """Open source."""
    if source.startswith("http"):
        webbrowser.open(source, new=2)
    elif source.endswith(".pdf"):
        webbrowser.open(f"file://{source}")
    else:
        webbrowser.open(f"file://{source}")


@cli.command()
def update_embeddings():
    """Update all the embeddings in your db.

    The only reason you would do this is if you change the embedding model, or
    the way the chunks are sized in your config.

    """
    config = get_config()
    db = get_db()
    from sentence_transformers import SentenceTransformer
    from langchain.text_splitter import RecursiveCharacterTextSplitter

    model = SentenceTransformer(config["embedding"]["model"])
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=config["embedding"]["chunk_size"],
        chunk_overlap=config["embedding"]["chunk_overlap"],
    )

    _, dim = model.encode(["test"]).shape

    # The point of this is to avoid deleting the database.
    db.execute("drop index if exists embedding_idx")
    db.execute("alter table sources drop embedding")
    db.execute(f"alter table sources add column embedding F32_BLOB({dim})")
    db.commit()

    for rowid, text in db.execute("select rowid, text from sources").fetchall():
        chunks = splitter.split_text(text)
        embedding = model.encode(chunks).mean(axis=0).astype(np.float32).tobytes()

        c = db.execute(
            "update sources set embedding = ? where rowid = ?", (embedding, rowid)
        )
        print(rowid, c.rowcount)

    # I don't know why this has to be here. I had it above, and no updates were
    # happening.
    db.execute(
        """create index if not exists embedding_idx ON sources (libsql_vector_idx(embedding))"""
    )
    db.commit()


######################
# Academic functions #
######################


@cli.command()
@click.argument("orcid")
def coa(orcid):
    """Generate Table 4 of Collaborators and Other Affiliations for NSF.

    ORCID is an orcid URL for the user to generate the table for.
    The file is saved in {orcid}-{today}.xlsx.
    """
    from .coa import get_coa

    get_coa(orcid)


@cli.command()
@click.argument("query", nargs=-1)
@click.option("-n", default=5, help="Number of documents to use")
def suggest_reviewers(query, n):
    """Suggest reviewers for QUERY.

    Use up to N similar documents. This is an iterative function, you will be
    prompted to expand the search.
    """
    config = get_config()
    query = " ".join(query)

    # This is a surprise. You can't just call the functions above! This is
    # apparently the way to do this.
    with click.Context(vsearch) as ctx:
        results = ctx.invoke(
            vsearch, query=query.split(" "), n=n, fmt="", iterative=True
        )

    # Now collect the authors from the matching papers
    authors = []

    for i, row in enumerate(results):
        source, citation, extra, distance = row

        d = json.loads(extra)

        for authorship in d["authorships"]:
            authors += [authorship["author"]["id"]]

    # get the unique ones
    from collections import Counter

    authors = Counter(authors)

    # Get author information
    data = []

    url = "https://api.openalex.org/authors/"
    # You can only filter on 50 ids at a time, so we hard code this limit here
    # and per page.
    for batch in batched(authors, 50):
        url = f'https://api.openalex.org/authors?filter=id:{"|".join(batch)}'

        params = {"per-page": 50, "email": config["openalex"]["email"]}

        r = get_data(url, params)

        for d in r["results"]:
            lki = d.get("last_known_institutions", [])
            if lki == []:
                affils = d.get("affiliations", [])
                if len(affils) >= 1:
                    lki = affils[0]["institution"]["display_name"]
                else:
                    lki = "unknown"

            else:
                if len(lki) >= 1:
                    lki = lki[0].get("display_name")
                else:
                    lki = "unknown"

            row = [
                d["display_name"],
                authors[d["id"]],
                d["summary_stats"]["h_index"],
                d["id"],
                lki,
            ]
            data += [row]

    # Sort and display the results
    data.sort(key=lambda row: row[2], reverse=True)

    if get_ipython():
        display(
            HTML(
                tabulate.tabulate(
                    data,
                    headers=["name", "# papers", "h-index", "oaid", "institution"],
                    tablefmt="html",
                )
            )
        )
        for i, row in enumerate(results):
            source, citation, extra, distance = row
            richprint(f"{i + 1:2d}. {citation} (source)\n\n")

    else:
        s = ["Potential reviewers"]
        s += [
            tabulate.tabulate(
                data,
                headers=["name", "# papers", "h-index", "oaid", "institution"],
                tablefmt="orgtbl",
            )
        ]
        s += ["\n" + "From these papers:"]
        for i, row in enumerate(results):
            source, citation, extra, distance = row
            s += [f"{i + 1:2d}. {citation} (source)\n\n"]

        console = Console(color_system="truecolor")
        with console.pager():
            for _s in s:
                console.print(_s)


@cli.command()
def app():
    """Launch the Streamlit app in the default web browser."""
    dirname = os.path.dirname(__file__)
    app = os.path.join(dirname, "app.py")
    os.system(f"streamlit run {app}")


@cli.command()
def version():
    """Print the version of litdb."""
    import pkg_resources

    version = pkg_resources.get_distribution("litdb").version
    print(f"Litdb: version {version}")


if __name__ == "__main__":
    cli()
